"""Merge Zenith proposal content into IDEALIZE Team Proposal template backgrounds."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "idealize" / "Team Proposal new (1).pptx"
OUTPUT = ROOT / "idealize" / "Zenith-Team-Proposal-IDEALIZE-2026.pptx"

NAVY = RGBColor(0x3B, 0x1C, 0x5A)
DARK = RGBColor(0x22, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED_NOTE = RGBColor(0xC0, 0x39, 0x2B)


def _fill(shape, text, *, size=11, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Calibri"
    run.font.color.rgb = color


def _box(slide, left, top, width, height, text, *, size=11, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(left, top, width, height)
    _fill(shape, text, size=size, bold=bold, color=color, align=align)
    return shape


def _member_block(n: int) -> str:
    return (
        f"Member {n}\n"
        f"Name:\n"
        f"University / School:\n"
        f"NIC No:\n"
        f"Email:\n"
        f"Contact Number:"
    )


PROBLEM_TEXT = (
    "Sri Lankan SMEs receive supplier invoices daily but struggle to time cheque payments "
    "against CBSL bank holidays and interbank clearing rules. Grouping invoices under LKR "
    "cheque ceilings is manual and error-prone. Cashiers need a simple mobile channel while "
    "finance teams require verified records before cheques are written."
)

BLUEPRINT_TEXT = (
    "Zenith is a Flask + SQLite web app with WhatsApp intake (Twilio). Cashiers send invoice "
    "photos; Gemini extracts fields; merchants verify on the dashboard. Python bundles "
    "invoices under an LKR ceiling, applies guardrails, and computes CBSL-aware liquidity "
    "dates. Features: invoice verification, cheque bundling chat, cash-flow deposit alerts, "
    "and analytics reports."
)

AGENTIC_TEXT = (
    "Five specialized agents: (1) Ingestion — Gemini vision reads invoices/cheques and supplier "
    "details; (2) Cheque Assistant — OpenAI bundling chat emits proposed_actions JSON; "
    "(3) Analyst — post-commit markdown reports; (4) Guide — in-app help only; "
    "(5) Liquidity Reviewer — audits bundles for max legal float. Python guardrails validate "
    "every action; humans verify before database commit."
)

TECH_TEXT = (
    "Stack: Python 3.12, Flask, SQLite, Google Gemini, OpenAI, Twilio WhatsApp. "
    "Core modules: guardrails.py, liquidity_engine.py (CBSL holidays), bundling.py, cash_flow.py. "
    "Feasibility: working prototype with web UI, WhatsApp webhook, multi-agent chat, and "
    "pending invoice verification flow."
)

MARKET_TEXT = (
    "Target users: Sri Lankan hardware merchants, finance staff, and field cashiers. "
    "Value: faster invoice intake via WhatsApp, smarter cheque dates, fewer surprise bank "
    "deposits, and extra working-capital float from holiday-aware scheduling."
)

SUSTAIN_TEXT = (
    "Subscription or per-merchant licensing for SMEs; expandable to more banks and suppliers. "
    "CBSL holiday sync keeps rules current. Agent prompts and guardrails can extend to payroll, "
    "supplier payments, and multi-branch operations without replacing human approval."
)


def build():
    prs = Presentation(str(TEMPLATE))

    # Slide 1 — cover (image6): subtitle under TEAM PROPOSAL bar
    s1 = prs.slides[0]
    _box(
        s1,
        Inches(1.2),
        Inches(3.05),
        Inches(7.6),
        Inches(0.55),
        "ZENITH — Agentic AI Cheque & Cash-Flow for Sri Lankan SMEs",
        size=16,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    _box(
        s1,
        Inches(1.2),
        Inches(3.55),
        Inches(7.6),
        Inches(0.35),
        "Team Name: [FILL IN]",
        size=13,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )

    # Slide 2 — proposal form table (image3)
    s2 = prs.slides[1]
    col_x = Inches(5.05)
    col_w = Inches(4.35)
    row_h = Inches(0.38)
    rows = [
        (Inches(1.88), "[FILL IN — your team name]"),
        (Inches(2.38), "Web App"),
        (Inches(2.88), "Open Category"),
        (Inches(3.38), "Zenith"),
        (Inches(3.88), "FinTech / SME Cash-Flow & Cheque Management"),
    ]
    for top, value in rows:
        _box(s2, col_x, top, col_w, row_h, value, size=11, color=DARK)

    # Slide 3 — Problem & Solution (image4)
    s3 = prs.slides[2]
    _box(s3, Inches(3.15), Inches(1.55), Inches(6.55), Inches(1.35), PROBLEM_TEXT, size=10)
    _box(s3, Inches(3.15), Inches(3.05), Inches(6.55), Inches(2.15), BLUEPRINT_TEXT, size=10)

    # Slide 4 — Innovation & Technical Design (image7)
    s4 = prs.slides[3]
    _box(s4, Inches(3.15), Inches(1.55), Inches(6.55), Inches(1.55), AGENTIC_TEXT, size=10)
    _box(s4, Inches(3.15), Inches(3.25), Inches(6.55), Inches(1.95), TECH_TEXT, size=10)

    # Slide 5 — Market Potential (image5)
    s5 = prs.slides[4]
    _box(s5, Inches(3.15), Inches(1.55), Inches(6.55), Inches(1.35), MARKET_TEXT, size=10)
    _box(s5, Inches(3.15), Inches(3.05), Inches(6.55), Inches(2.0), SUSTAIN_TEXT, size=10)

    # Slides 6–8 — team member blocks (existing text boxes)
    member_idx = 0
    for slide in (prs.slides[5], prs.slides[6], prs.slides[7]):
        for shape in slide.shapes:
            if shape.has_text_frame and "Team Leader Name" in shape.text:
                member_idx += 1
                _fill(shape, _member_block(member_idx), size=11, color=DARK)

    # Slide 9 — optional extra (image8): one-line pitch
    s9 = prs.slides[8]
    _box(
        s9,
        Inches(1.0),
        Inches(2.2),
        Inches(8.0),
        Inches(2.0),
        "Pitch: Zenith is a multi-agent SME finance copilot that reads invoices on WhatsApp, "
        "bundles cheques under guardrails, and maximizes legal cash float using CBSL clearing rules.",
        size=14,
        bold=True,
        color=NAVY,
        align=PP_ALIGN.CENTER,
    )

    # Slide 10 — Thank You (image2): leave template as-is

    prs.save(str(OUTPUT))
    print(f"Saved merged proposal: {OUTPUT}")


if __name__ == "__main__":
    build()
